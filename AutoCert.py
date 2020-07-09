import argparse
import csv
from pptx import Presentation

parser = argparse.ArgumentParser(description='Automatic certiicate generator')
parser.add_argument('t', metavar='template', type=str)
parser.add_argument('n',metavar='names', type=str)
parser.add_argument('-d',metavar='placeholder', type=str, default = "Jane Smith")
parser.add_argument('-dem',metavar='delimiter', type=str, default = '\n')
args = parser.parse_args()

def search_and_replace(search_str, repl_str, input, output):
    """
    Replaces all occurence of search_str with repl_str in the input pptx file
    and saves it with output as the name.
    """
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(search_str), str(repl_str))
                            run.text = new_text
    prs.save(output)


with open(args.n) as csv_file:
    csv_reader = csv.reader(csv_file, delimiter=args.dem)
    for name in csv_reader:
        search_and_replace(args.d, name[0] , args.t ,name[0]+".pptx")